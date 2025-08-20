"""
Executive Reporting Module for Patent #3: Unified AI-Driven Cloud Governance Platform

This module provides automated executive-level reporting with customizable templates,
compliance attestation, board-ready visualizations, and scheduled delivery.

Key Features:
- Automated report generation
- Customizable report templates
- Scheduled report delivery
- Compliance attestation reports
- Board-ready visualizations
- Multi-format export (PDF, PowerPoint, Excel)
"""

import asyncio
import json
import logging
from datetime import datetime, timedelta
from typing import Dict, List, Any, Optional, Tuple
from enum import Enum
from dataclasses import dataclass, field
import numpy as np
import pandas as pd
from pathlib import Path
import base64
import io

# Visualization libraries
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.figure import Figure
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots

# Document generation
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4, landscape
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.pdfgen import canvas
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.linecharts import HorizontalLineChart

# Email and scheduling
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import schedule
from jinja2 import Template

# Export to PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Export to Excel
import openpyxl
from openpyxl.styles import Font, Fill, PatternFill, Border, Side, Alignment
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.utils.dataframe import dataframe_to_rows

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportType(Enum):
    """Types of executive reports"""
    EXECUTIVE_SUMMARY = "executive_summary"
    COMPLIANCE_ATTESTATION = "compliance_attestation"
    BOARD_PRESENTATION = "board_presentation"
    QUARTERLY_REVIEW = "quarterly_review"
    RISK_ASSESSMENT = "risk_assessment"
    COST_OPTIMIZATION = "cost_optimization"
    SECURITY_POSTURE = "security_posture"
    OPERATIONAL_EXCELLENCE = "operational_excellence"
    AUDIT_REPORT = "audit_report"
    GOVERNANCE_SCORECARD = "governance_scorecard"


class ReportFormat(Enum):
    """Output formats for reports"""
    PDF = "pdf"
    POWERPOINT = "pptx"
    EXCEL = "xlsx"
    HTML = "html"
    MARKDOWN = "md"
    JSON = "json"


class DeliveryMethod(Enum):
    """Report delivery methods"""
    EMAIL = "email"
    TEAMS = "teams"
    SLACK = "slack"
    SHAREPOINT = "sharepoint"
    S3 = "s3"
    AZURE_BLOB = "azure_blob"
    API_WEBHOOK = "webhook"


@dataclass
class ReportTemplate:
    """Customizable report template"""
    name: str
    type: ReportType
    sections: List[str]
    metrics: List[str]
    visualizations: List[str]
    branding: Dict[str, Any] = field(default_factory=dict)
    custom_css: Optional[str] = None
    custom_headers: Optional[Dict[str, str]] = None


@dataclass
class ReportSchedule:
    """Report scheduling configuration"""
    report_type: ReportType
    frequency: str  # daily, weekly, monthly, quarterly
    time: str  # HH:MM format
    recipients: List[str]
    format: ReportFormat
    delivery_method: DeliveryMethod
    filters: Optional[Dict[str, Any]] = None
    template: Optional[ReportTemplate] = None


@dataclass
class ComplianceAttestation:
    """Compliance attestation details"""
    framework: str
    compliance_level: float
    attestation_date: datetime
    attested_by: str
    findings: List[Dict[str, Any]]
    remediation_plan: Optional[Dict[str, Any]] = None
    evidence_links: List[str] = field(default_factory=list)


class ExecutiveReportingEngine:
    """
    Main executive reporting engine for generating and distributing
    governance reports to leadership and board members.
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the executive reporting engine"""
        self.config = config or {}
        self.templates: Dict[str, ReportTemplate] = {}
        self.schedules: List[ReportSchedule] = []
        self.report_cache: Dict[str, Any] = {}
        
        # Initialize components
        self._init_templates()
        self._init_visualization_engine()
        self._init_delivery_system()
        
    def _init_templates(self):
        """Initialize default report templates"""
        # Executive Summary Template
        self.templates['executive_summary'] = ReportTemplate(
            name="Executive Summary",
            type=ReportType.EXECUTIVE_SUMMARY,
            sections=[
                "Key Metrics Overview",
                "Risk Highlights",
                "Compliance Status",
                "Cost Analysis",
                "Recommendations",
                "Trend Analysis"
            ],
            metrics=[
                "overall_governance_score",
                "security_posture",
                "compliance_percentage",
                "monthly_spend",
                "critical_risks",
                "open_incidents"
            ],
            visualizations=[
                "governance_scorecard",
                "risk_heatmap",
                "compliance_radar",
                "cost_trend",
                "incident_timeline"
            ],
            branding={
                "logo": "assets/logo.png",
                "primary_color": "#003366",
                "secondary_color": "#0066CC",
                "font_family": "Helvetica"
            }
        )
        
        # Board Presentation Template
        self.templates['board_presentation'] = ReportTemplate(
            name="Board Presentation",
            type=ReportType.BOARD_PRESENTATION,
            sections=[
                "Executive Overview",
                "Strategic Initiatives",
                "Risk Management",
                "Compliance & Regulatory",
                "Financial Impact",
                "Future Outlook"
            ],
            metrics=[
                "strategic_kpis",
                "risk_exposure",
                "regulatory_compliance",
                "cost_savings",
                "efficiency_gains",
                "maturity_score"
            ],
            visualizations=[
                "strategic_dashboard",
                "risk_matrix",
                "compliance_timeline",
                "roi_analysis",
                "maturity_progression"
            ],
            branding={
                "logo": "assets/logo.png",
                "theme": "corporate",
                "slide_layout": "professional"
            }
        )
        
        # Compliance Attestation Template
        self.templates['compliance_attestation'] = ReportTemplate(
            name="Compliance Attestation",
            type=ReportType.COMPLIANCE_ATTESTATION,
            sections=[
                "Attestation Statement",
                "Compliance Overview",
                "Control Assessment",
                "Findings & Exceptions",
                "Remediation Status",
                "Evidence Documentation"
            ],
            metrics=[
                "control_effectiveness",
                "compliance_score",
                "exception_count",
                "remediation_progress",
                "audit_findings"
            ],
            visualizations=[
                "control_matrix",
                "compliance_trend",
                "exception_chart",
                "remediation_timeline"
            ]
        )
    
    def _init_visualization_engine(self):
        """Initialize visualization engine with custom styles"""
        # Set default styles
        sns.set_theme(style="whitegrid")
        plt.rcParams['figure.figsize'] = (12, 6)
        plt.rcParams['font.size'] = 10
        
        # Custom color palettes
        self.color_palettes = {
            'governance': ['#003366', '#0066CC', '#3399FF', '#66B2FF', '#99CCFF'],
            'risk': ['#00FF00', '#FFFF00', '#FFA500', '#FF0000', '#8B0000'],
            'compliance': ['#28a745', '#ffc107', '#dc3545'],
            'cost': ['#17a2b8', '#20c997', '#6610f2', '#e83e8c']
        }
    
    def _init_delivery_system(self):
        """Initialize report delivery system"""
        self.delivery_config = {
            'email': {
                'smtp_server': self.config.get('smtp_server', 'smtp.gmail.com'),
                'smtp_port': self.config.get('smtp_port', 587),
                'sender': self.config.get('sender_email'),
                'password': self.config.get('email_password')
            },
            'teams': {
                'webhook_url': self.config.get('teams_webhook')
            },
            'slack': {
                'webhook_url': self.config.get('slack_webhook')
            }
        }
    
    async def generate_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any],
        format: ReportFormat = ReportFormat.PDF,
        template: Optional[ReportTemplate] = None
    ) -> bytes:
        """
        Generate an executive report based on type and data
        
        Args:
            report_type: Type of report to generate
            data: Report data including metrics and insights
            format: Output format for the report
            template: Optional custom template
            
        Returns:
            Generated report as bytes
        """
        logger.info(f"Generating {report_type.value} report in {format.value} format")
        
        # Use default template if not provided
        if template is None:
            template = self.templates.get(report_type.value)
        
        # Generate report based on format
        if format == ReportFormat.PDF:
            return await self._generate_pdf_report(report_type, data, template)
        elif format == ReportFormat.POWERPOINT:
            return await self._generate_pptx_report(report_type, data, template)
        elif format == ReportFormat.EXCEL:
            return await self._generate_excel_report(report_type, data, template)
        elif format == ReportFormat.HTML:
            return await self._generate_html_report(report_type, data, template)
        else:
            return await self._generate_json_report(report_type, data)
    
    async def _generate_pdf_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any],
        template: ReportTemplate
    ) -> bytes:
        """Generate PDF report"""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor(template.branding.get('primary_color', '#003366')),
            alignment=TA_CENTER,
            spaceAfter=30
        )
        
        # Add title
        title = Paragraph(template.name, title_style)
        story.append(title)
        story.append(Spacer(1, 0.3*inch))
        
        # Add generation timestamp
        timestamp = Paragraph(
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            styles['Normal']
        )
        story.append(timestamp)
        story.append(Spacer(1, 0.5*inch))
        
        # Process each section
        for section in template.sections:
            # Section header
            section_header = Paragraph(section, styles['Heading2'])
            story.append(section_header)
            story.append(Spacer(1, 0.2*inch))
            
            # Add section content based on type
            if section == "Key Metrics Overview":
                metrics_table = self._create_metrics_table(data.get('metrics', {}))
                story.append(metrics_table)
            elif section == "Risk Highlights":
                risk_chart = self._create_risk_visualization(data.get('risks', []))
                story.append(risk_chart)
            elif section == "Compliance Status":
                compliance_data = self._create_compliance_summary(data.get('compliance', {}))
                story.append(compliance_data)
            elif section == "Cost Analysis":
                cost_chart = self._create_cost_analysis(data.get('costs', {}))
                story.append(cost_chart)
            elif section == "Recommendations":
                recommendations = self._format_recommendations(data.get('recommendations', []))
                story.append(recommendations)
            
            story.append(PageBreak())
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    
    async def _generate_pptx_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any],
        template: ReportTemplate
    ) -> bytes:
        """Generate PowerPoint presentation"""
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = template.name
        subtitle.text = f"Governance Report - {datetime.now().strftime('%B %Y')}"
        
        # Process each section as a slide
        for section in template.sections:
            # Add content slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            # Set slide title
            title_shape = shapes.title
            title_shape.text = section
            
            # Add content based on section
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            
            if section == "Executive Overview":
                self._add_executive_overview_slide(tf, data)
            elif section == "Risk Management":
                self._add_risk_management_slide(slide, data)
            elif section == "Compliance & Regulatory":
                self._add_compliance_slide(slide, data)
            elif section == "Financial Impact":
                self._add_financial_slide(slide, data)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    async def _generate_excel_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any],
        template: ReportTemplate
    ) -> bytes:
        """Generate Excel report with multiple sheets"""
        wb = openpyxl.Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Executive Summary Sheet
        ws_summary = wb.create_sheet("Executive Summary")
        self._create_summary_sheet(ws_summary, data)
        
        # Metrics Sheet
        ws_metrics = wb.create_sheet("Metrics")
        self._create_metrics_sheet(ws_metrics, data.get('metrics', {}))
        
        # Risks Sheet
        ws_risks = wb.create_sheet("Risks")
        self._create_risks_sheet(ws_risks, data.get('risks', []))
        
        # Compliance Sheet
        ws_compliance = wb.create_sheet("Compliance")
        self._create_compliance_sheet(ws_compliance, data.get('compliance', {}))
        
        # Recommendations Sheet
        ws_recommendations = wb.create_sheet("Recommendations")
        self._create_recommendations_sheet(ws_recommendations, data.get('recommendations', []))
        
        # Save to bytes
        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    async def _generate_html_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any],
        template: ReportTemplate
    ) -> bytes:
        """Generate HTML report with interactive visualizations"""
        html_template = Template("""
        <!DOCTYPE html>
        <html>
        <head>
            <title>{{ title }}</title>
            <meta charset="utf-8">
            <style>
                body { font-family: {{ font_family }}; margin: 40px; }
                h1 { color: {{ primary_color }}; }
                h2 { color: {{ secondary_color }}; }
                .metric-card {
                    display: inline-block;
                    padding: 20px;
                    margin: 10px;
                    border: 1px solid #ddd;
                    border-radius: 8px;
                    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                }
                .metric-value { font-size: 2em; font-weight: bold; }
                .metric-label { color: #666; }
                .chart-container { margin: 30px 0; }
                table { border-collapse: collapse; width: 100%; }
                th, td { border: 1px solid #ddd; padding: 12px; text-align: left; }
                th { background-color: {{ primary_color }}; color: white; }
                .risk-high { background-color: #ffebee; }
                .risk-medium { background-color: #fff3e0; }
                .risk-low { background-color: #e8f5e9; }
            </style>
            <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
        </head>
        <body>
            <h1>{{ title }}</h1>
            <p>Generated: {{ timestamp }}</p>
            
            {% for section in sections %}
            <section>
                <h2>{{ section.title }}</h2>
                {{ section.content | safe }}
            </section>
            {% endfor %}
            
            <footer>
                <p>&copy; 2024 PolicyCortex. All rights reserved.</p>
            </footer>
        </body>
        </html>
        """)
        
        # Prepare sections
        sections = []
        for section_name in template.sections:
            section_content = await self._generate_html_section(section_name, data)
            sections.append({
                'title': section_name,
                'content': section_content
            })
        
        # Render HTML
        html = html_template.render(
            title=template.name,
            font_family=template.branding.get('font_family', 'Arial'),
            primary_color=template.branding.get('primary_color', '#003366'),
            secondary_color=template.branding.get('secondary_color', '#0066CC'),
            timestamp=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            sections=sections
        )
        
        return html.encode('utf-8')
    
    async def _generate_json_report(
        self,
        report_type: ReportType,
        data: Dict[str, Any]
    ) -> bytes:
        """Generate JSON report for API consumption"""
        report = {
            'report_type': report_type.value,
            'generated_at': datetime.now().isoformat(),
            'data': data,
            'metadata': {
                'version': '1.0',
                'format': 'json',
                'schema': 'executive_report_v1'
            }
        }
        return json.dumps(report, indent=2).encode('utf-8')
    
    def _create_metrics_table(self, metrics: Dict[str, Any]) -> Table:
        """Create metrics table for PDF report"""
        data = [['Metric', 'Value', 'Trend', 'Status']]
        
        for metric_name, metric_data in metrics.items():
            trend_symbol = '↑' if metric_data.get('trend') == 'up' else '↓'
            status = '✓' if metric_data.get('status') == 'good' else '⚠'
            data.append([
                metric_name.replace('_', ' ').title(),
                str(metric_data.get('value', 'N/A')),
                trend_symbol,
                status
            ])
        
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        return table
    
    def _create_risk_visualization(self, risks: List[Dict[str, Any]]) -> Drawing:
        """Create risk heatmap visualization"""
        drawing = Drawing(400, 200)
        
        # Create risk matrix
        risk_matrix = [[0] * 5 for _ in range(5)]
        
        for risk in risks:
            likelihood = risk.get('likelihood', 1) - 1
            impact = risk.get('impact', 1) - 1
            risk_matrix[likelihood][impact] += 1
        
        # Create heatmap chart
        # (Simplified for this implementation)
        return drawing
    
    def _create_compliance_summary(self, compliance: Dict[str, Any]) -> Table:
        """Create compliance summary table"""
        data = [['Framework', 'Compliance %', 'Controls', 'Exceptions']]
        
        for framework, details in compliance.items():
            data.append([
                framework,
                f"{details.get('percentage', 0):.1f}%",
                str(details.get('controls', 0)),
                str(details.get('exceptions', 0))
            ])
        
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        return table
    
    def _create_cost_analysis(self, costs: Dict[str, Any]) -> Drawing:
        """Create cost analysis chart"""
        drawing = Drawing(400, 200)
        
        # Create bar chart
        bc = VerticalBarChart()
        bc.x = 50
        bc.y = 50
        bc.height = 125
        bc.width = 300
        
        # Add data
        data = []
        categories = []
        for category, amount in costs.items():
            categories.append(category)
            data.append([amount])
        
        bc.data = data
        bc.categoryAxis.categoryNames = categories
        
        drawing.add(bc)
        return drawing
    
    def _format_recommendations(self, recommendations: List[Dict[str, Any]]) -> Table:
        """Format recommendations as table"""
        data = [['Priority', 'Recommendation', 'Impact', 'Effort']]
        
        for rec in recommendations[:10]:  # Top 10 recommendations
            data.append([
                rec.get('priority', 'Medium'),
                rec.get('title', ''),
                rec.get('impact', 'Medium'),
                rec.get('effort', 'Medium')
            ])
        
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        return table
    
    async def schedule_report(
        self,
        schedule: ReportSchedule
    ) -> None:
        """
        Schedule a report for automatic generation and delivery
        
        Args:
            schedule: Report scheduling configuration
        """
        self.schedules.append(schedule)
        
        # Schedule based on frequency
        if schedule.frequency == 'daily':
            schedule.every().day.at(schedule.time).do(
                lambda: asyncio.create_task(self._execute_scheduled_report(schedule))
            )
        elif schedule.frequency == 'weekly':
            schedule.every().monday.at(schedule.time).do(
                lambda: asyncio.create_task(self._execute_scheduled_report(schedule))
            )
        elif schedule.frequency == 'monthly':
            schedule.every().month.do(
                lambda: asyncio.create_task(self._execute_scheduled_report(schedule))
            )
        
        logger.info(f"Scheduled {schedule.report_type.value} report for {schedule.frequency} delivery")
    
    async def _execute_scheduled_report(self, schedule: ReportSchedule) -> None:
        """Execute a scheduled report"""
        try:
            # Gather data for report
            data = await self._gather_report_data(schedule.report_type, schedule.filters)
            
            # Generate report
            report = await self.generate_report(
                schedule.report_type,
                data,
                schedule.format,
                schedule.template
            )
            
            # Deliver report
            await self.deliver_report(
                report,
                schedule.recipients,
                schedule.delivery_method,
                schedule.report_type
            )
            
            logger.info(f"Successfully executed scheduled report: {schedule.report_type.value}")
            
        except Exception as e:
            logger.error(f"Failed to execute scheduled report: {str(e)}")
    
    async def deliver_report(
        self,
        report: bytes,
        recipients: List[str],
        method: DeliveryMethod,
        report_type: ReportType
    ) -> None:
        """
        Deliver report to recipients via specified method
        
        Args:
            report: Report content as bytes
            recipients: List of recipient addresses
            method: Delivery method
            report_type: Type of report for naming
        """
        if method == DeliveryMethod.EMAIL:
            await self._deliver_via_email(report, recipients, report_type)
        elif method == DeliveryMethod.TEAMS:
            await self._deliver_via_teams(report, report_type)
        elif method == DeliveryMethod.SLACK:
            await self._deliver_via_slack(report, report_type)
        elif method == DeliveryMethod.SHAREPOINT:
            await self._deliver_via_sharepoint(report, report_type)
        else:
            logger.warning(f"Delivery method {method.value} not implemented")
    
    async def _deliver_via_email(
        self,
        report: bytes,
        recipients: List[str],
        report_type: ReportType
    ) -> None:
        """Deliver report via email"""
        msg = MIMEMultipart()
        msg['From'] = self.delivery_config['email']['sender']
        msg['To'] = ', '.join(recipients)
        msg['Subject'] = f"{report_type.value.replace('_', ' ').title()} - {datetime.now().strftime('%B %Y')}"
        
        # Add body
        body = f"""
        Please find attached the {report_type.value.replace('_', ' ').title()} report.
        
        Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
        
        This is an automated report from PolicyCortex Executive Reporting System.
        """
        msg.attach(MIMEText(body, 'plain'))
        
        # Add attachment
        attachment = MIMEBase('application', 'octet-stream')
        attachment.set_payload(report)
        encoders.encode_base64(attachment)
        attachment.add_header(
            'Content-Disposition',
            f'attachment; filename={report_type.value}_{datetime.now().strftime("%Y%m%d")}.pdf'
        )
        msg.attach(attachment)
        
        # Send email
        server = smtplib.SMTP(
            self.delivery_config['email']['smtp_server'],
            self.delivery_config['email']['smtp_port']
        )
        server.starttls()
        server.login(
            self.delivery_config['email']['sender'],
            self.delivery_config['email']['password']
        )
        server.send_message(msg)
        server.quit()
        
        logger.info(f"Report delivered via email to {len(recipients)} recipients")
    
    async def generate_compliance_attestation(
        self,
        framework: str,
        assessment_data: Dict[str, Any],
        attester: str
    ) -> ComplianceAttestation:
        """
        Generate a formal compliance attestation report
        
        Args:
            framework: Compliance framework (SOC2, ISO27001, etc.)
            assessment_data: Assessment results and findings
            attester: Name of person attesting
            
        Returns:
            Compliance attestation object
        """
        attestation = ComplianceAttestation(
            framework=framework,
            compliance_level=assessment_data.get('compliance_score', 0),
            attestation_date=datetime.now(),
            attested_by=attester,
            findings=assessment_data.get('findings', []),
            remediation_plan=assessment_data.get('remediation_plan'),
            evidence_links=assessment_data.get('evidence_links', [])
        )
        
        # Generate formal attestation document
        attestation_doc = await self._generate_attestation_document(attestation)
        
        # Store attestation for audit trail
        await self._store_attestation(attestation, attestation_doc)
        
        return attestation
    
    async def _generate_attestation_document(
        self,
        attestation: ComplianceAttestation
    ) -> bytes:
        """Generate formal attestation document"""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title = Paragraph(
            f"{attestation.framework} Compliance Attestation",
            styles['Title']
        )
        story.append(title)
        story.append(Spacer(1, 0.5*inch))
        
        # Attestation statement
        statement = f"""
        I, {attestation.attested_by}, hereby attest that the PolicyCortex platform
        has been assessed for compliance with {attestation.framework} requirements.
        
        Based on the assessment conducted on {attestation.attestation_date.strftime('%Y-%m-%d')},
        the platform demonstrates a compliance level of {attestation.compliance_level:.1f}%.
        
        This attestation is based on the evidence collected and controls tested
        as documented in the accompanying assessment report.
        """
        
        story.append(Paragraph(statement, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Findings summary
        if attestation.findings:
            story.append(Paragraph("Key Findings:", styles['Heading2']))
            for finding in attestation.findings[:5]:  # Top 5 findings
                story.append(Paragraph(f"• {finding.get('description', '')}", styles['Normal']))
        
        # Signature block
        story.append(Spacer(1, 1*inch))
        story.append(Paragraph("_" * 40, styles['Normal']))
        story.append(Paragraph(attestation.attested_by, styles['Normal']))
        story.append(Paragraph(
            attestation.attestation_date.strftime('%Y-%m-%d'),
            styles['Normal']
        ))
        
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    
    async def _store_attestation(
        self,
        attestation: ComplianceAttestation,
        document: bytes
    ) -> None:
        """Store attestation for audit trail"""
        # Store in database or document management system
        # Implementation depends on backend infrastructure
        pass
    
    async def _gather_report_data(
        self,
        report_type: ReportType,
        filters: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """Gather data for report generation"""
        # This would integrate with various data sources
        # For now, return sample data
        return {
            'metrics': {
                'governance_score': {'value': 85.2, 'trend': 'up', 'status': 'good'},
                'security_posture': {'value': 92.1, 'trend': 'up', 'status': 'good'},
                'compliance_percentage': {'value': 94.5, 'trend': 'stable', 'status': 'good'},
                'monthly_spend': {'value': 125000, 'trend': 'down', 'status': 'good'}
            },
            'risks': [
                {'name': 'Data Breach', 'likelihood': 2, 'impact': 5, 'status': 'mitigated'},
                {'name': 'Compliance Violation', 'likelihood': 1, 'impact': 4, 'status': 'monitoring'}
            ],
            'compliance': {
                'SOC2': {'percentage': 95.0, 'controls': 120, 'exceptions': 6},
                'ISO27001': {'percentage': 92.0, 'controls': 114, 'exceptions': 9},
                'GDPR': {'percentage': 98.0, 'controls': 45, 'exceptions': 1}
            },
            'costs': {
                'Compute': 45000,
                'Storage': 23000,
                'Network': 12000,
                'Security': 28000,
                'Other': 17000
            },
            'recommendations': [
                {
                    'priority': 'High',
                    'title': 'Enable MFA for all admin accounts',
                    'impact': 'High',
                    'effort': 'Low'
                },
                {
                    'priority': 'Medium',
                    'title': 'Optimize storage lifecycle policies',
                    'impact': 'Medium',
                    'effort': 'Medium'
                }
            ]
        }
    
    async def generate_board_deck(
        self,
        quarter: str,
        year: int,
        data: Dict[str, Any]
    ) -> bytes:
        """
        Generate quarterly board presentation deck
        
        Args:
            quarter: Quarter (Q1, Q2, Q3, Q4)
            year: Year
            data: Quarterly data and metrics
            
        Returns:
            PowerPoint presentation as bytes
        """
        prs = Presentation()
        
        # Title slide
        self._add_title_slide(prs, f"{quarter} {year} Governance Review")
        
        # Executive summary
        self._add_executive_summary_slide(prs, data)
        
        # Key metrics
        self._add_metrics_slide(prs, data.get('metrics', {}))
        
        # Risk overview
        self._add_risk_overview_slide(prs, data.get('risks', []))
        
        # Compliance status
        self._add_compliance_status_slide(prs, data.get('compliance', {}))
        
        # Cost optimization
        self._add_cost_optimization_slide(prs, data.get('costs', {}))
        
        # Strategic initiatives
        self._add_strategic_initiatives_slide(prs, data.get('initiatives', []))
        
        # Q&A slide
        self._add_qa_slide(prs)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    def _add_title_slide(self, prs: Presentation, title_text: str):
        """Add title slide to presentation"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = f"PolicyCortex Executive Report\n{datetime.now().strftime('%B %d, %Y')}"
    
    def _add_executive_summary_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Add executive summary slide"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = "Executive Summary"
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        # Add key points
        points = [
            f"Overall Governance Score: {data.get('governance_score', 0):.1f}%",
            f"Critical Risks Identified: {data.get('critical_risks', 0)}",
            f"Compliance Level: {data.get('compliance_level', 0):.1f}%",
            f"Cost Optimization Potential: ${data.get('savings_potential', 0):,.0f}",
            f"Security Incidents: {data.get('security_incidents', 0)}"
        ]
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    def _add_metrics_slide(self, prs: Presentation, metrics: Dict[str, Any]):
        """Add key metrics slide"""
        # Would add charts and visualizations
        pass
    
    def _add_risk_overview_slide(self, prs: Presentation, risks: List[Dict[str, Any]]):
        """Add risk overview slide"""
        # Would add risk matrix visualization
        pass
    
    def _add_compliance_status_slide(self, prs: Presentation, compliance: Dict[str, Any]):
        """Add compliance status slide"""
        # Would add compliance charts
        pass
    
    def _add_cost_optimization_slide(self, prs: Presentation, costs: Dict[str, Any]):
        """Add cost optimization slide"""
        # Would add cost breakdown charts
        pass
    
    def _add_strategic_initiatives_slide(self, prs: Presentation, initiatives: List[Dict[str, Any]]):
        """Add strategic initiatives slide"""
        # Would add initiatives timeline
        pass
    
    def _add_qa_slide(self, prs: Presentation):
        """Add Q&A slide"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Questions & Discussion"
    
    def _create_summary_sheet(self, ws, data: Dict[str, Any]):
        """Create executive summary Excel sheet"""
        ws['A1'] = 'Executive Summary'
        ws['A1'].font = Font(size=16, bold=True)
        
        # Add metrics
        row = 3
        ws[f'A{row}'] = 'Key Metrics'
        ws[f'A{row}'].font = Font(bold=True)
        row += 1
        
        for metric, value in data.get('metrics', {}).items():
            ws[f'A{row}'] = metric.replace('_', ' ').title()
            ws[f'B{row}'] = value.get('value', 'N/A')
            row += 1
    
    def _create_metrics_sheet(self, ws, metrics: Dict[str, Any]):
        """Create detailed metrics Excel sheet"""
        # Add headers
        headers = ['Metric', 'Value', 'Target', 'Status', 'Trend']
        for col, header in enumerate(headers, 1):
            ws.cell(row=1, column=col, value=header)
            ws.cell(row=1, column=col).font = Font(bold=True)
            ws.cell(row=1, column=col).fill = PatternFill(
                start_color='366092',
                end_color='366092',
                fill_type='solid'
            )
    
    def _create_risks_sheet(self, ws, risks: List[Dict[str, Any]]):
        """Create risks Excel sheet"""
        # Add risk data
        pass
    
    def _create_compliance_sheet(self, ws, compliance: Dict[str, Any]):
        """Create compliance Excel sheet"""
        # Add compliance data
        pass
    
    def _create_recommendations_sheet(self, ws, recommendations: List[Dict[str, Any]]):
        """Create recommendations Excel sheet"""
        # Add recommendations
        pass
    
    async def _generate_html_section(self, section_name: str, data: Dict[str, Any]) -> str:
        """Generate HTML content for a report section"""
        if section_name == "Key Metrics Overview":
            return self._generate_metrics_html(data.get('metrics', {}))
        elif section_name == "Risk Highlights":
            return self._generate_risks_html(data.get('risks', []))
        elif section_name == "Compliance Status":
            return self._generate_compliance_html(data.get('compliance', {}))
        else:
            return f"<p>Section {section_name} content</p>"
    
    def _generate_metrics_html(self, metrics: Dict[str, Any]) -> str:
        """Generate HTML for metrics section"""
        html = '<div class="metrics-container">'
        for metric, value in metrics.items():
            html += f'''
            <div class="metric-card">
                <div class="metric-label">{metric.replace('_', ' ').title()}</div>
                <div class="metric-value">{value.get('value', 'N/A')}</div>
            </div>
            '''
        html += '</div>'
        return html
    
    def _generate_risks_html(self, risks: List[Dict[str, Any]]) -> str:
        """Generate HTML for risks section"""
        html = '<table><thead><tr><th>Risk</th><th>Likelihood</th><th>Impact</th><th>Status</th></tr></thead><tbody>'
        for risk in risks:
            risk_class = 'risk-high' if risk.get('impact', 0) >= 4 else 'risk-medium' if risk.get('impact', 0) >= 2 else 'risk-low'
            html += f'''
            <tr class="{risk_class}">
                <td>{risk.get('name', '')}</td>
                <td>{risk.get('likelihood', '')}</td>
                <td>{risk.get('impact', '')}</td>
                <td>{risk.get('status', '')}</td>
            </tr>
            '''
        html += '</tbody></table>'
        return html
    
    def _generate_compliance_html(self, compliance: Dict[str, Any]) -> str:
        """Generate HTML for compliance section"""
        html = '<table><thead><tr><th>Framework</th><th>Compliance %</th><th>Controls</th><th>Exceptions</th></tr></thead><tbody>'
        for framework, details in compliance.items():
            html += f'''
            <tr>
                <td>{framework}</td>
                <td>{details.get('percentage', 0):.1f}%</td>
                <td>{details.get('controls', 0)}</td>
                <td>{details.get('exceptions', 0)}</td>
            </tr>
            '''
        html += '</tbody></table>'
        return html


# Export main class
__all__ = ['ExecutiveReportingEngine', 'ReportType', 'ReportFormat', 'DeliveryMethod']